import streamlit as st
from dotenv import load_dotenv
# from langchain_groq import ChatGroq
from langchain_community.llms import HuggingFacePipeline
from transformers import pipeline

# Initialize HF pipeline
pipe = pipeline(
    "text-generation",
    model="TheBloke/Llama-3-1-HF",
    max_new_tokens=512,
    temperature=0.2
)
llm = HuggingFacePipeline(pipeline=pipe)

from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.output_parsers import PydanticOutputParser
from sklearn.metrics.pairwise import cosine_similarity
from pydantic import BaseModel, Field
from typing import List
from pypdf import PdfReader
import pptx

# ----------------------------------------------------
# PAGE CONFIG
# ----------------------------------------------------
st.set_page_config(
    page_title="AI Document Search",
    page_icon="📘",
    layout="wide"
)

# ----------------------------------------------------
# LOAD ENV
# ----------------------------------------------------
load_dotenv()

# ----------------------------------------------------
# DATA MODEL
# ----------------------------------------------------
class RAGAnswer(BaseModel):
    answer: str = Field(description="Answer to the question")
    confidence: str = Field(description="high / medium / low")
    source_chunks: List[str] = Field(description="Top relevant chunks")

# ----------------------------------------------------
# LOAD MODELS
# ----------------------------------------------------
@st.cache_resource
def load_models():

    # ---------- LLM ----------
    llm = ChatGroq(
        model="llama-3.1-8b-instant",
        temperature=0.2,
        max_tokens=512  # limit response tokens
    )

    parser = PydanticOutputParser(pydantic_object=RAGAnswer)

    # ---------- Embeddings ----------
    primary_embeddings = GoogleGenerativeAIEmbeddings(
        model="gemini-embedding-001"
    )

    fallback_embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        model_kwargs={"device": "cpu"}  # use "cuda" if GPU available
    )

    return llm, primary_embeddings, fallback_embeddings, parser


llm, primary_embeddings, fallback_embeddings, parser = load_models()

# ----------------------------------------------------
# SILENT EMBEDDING FALLBACK
# ----------------------------------------------------
def safe_embed_documents(texts):
    try:
        return primary_embeddings.embed_documents(texts)
    except Exception:
        return fallback_embeddings.embed_documents(texts)


def safe_embed_query(query):
    try:
        return primary_embeddings.embed_query(query)
    except Exception:
        return fallback_embeddings.embed_query(query)

# ----------------------------------------------------
# TEXT CHUNKING
# ----------------------------------------------------
def chunk_text(text, chunk_size=500, overlap=100):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return chunks


def process_documents(text):
    return chunk_text(text, chunk_size=500, overlap=100)

# ----------------------------------------------------
# RAG QUERY (TOKEN SAFE)
# ----------------------------------------------------
def rag_query(chunks, query, top_k=3):

    doc_embeddings = safe_embed_documents(chunks)
    query_embedding = safe_embed_query(query)

    scores = cosine_similarity([query_embedding], doc_embeddings)[0]
    top_indices = scores.argsort()[-top_k:][::-1]

    sources = [chunks[i] for i in top_indices]

    # limit total context to prevent Groq token limit errors
    MAX_CONTEXT_CHARS = 2000
    context = ""
    for src in sources:
        if len(context) + len(src) > MAX_CONTEXT_CHARS:
            break
        context += src + "\n\n"

    prompt = f"""
Answer ONLY from the context.

CONTEXT:
{context}

QUESTION:
{query}

{parser.get_format_instructions()}
"""

    chain = llm | parser
    return chain.invoke(prompt)

# ----------------------------------------------------
# UI HEADER
# ----------------------------------------------------
st.title("📘 AI Document Search System")
st.caption("Groq Llama-3.1 | Token-safe RAG with silent embedding fallback")

# ----------------------------------------------------
# FILE UPLOAD
# ----------------------------------------------------
st.sidebar.header("📂 Upload Documents")

uploaded_files = st.sidebar.file_uploader(
    "Upload TXT / PDF / PPT",
    type=["txt", "pdf", "pptx"],
    accept_multiple_files=True
)

if uploaded_files:
    full_text = ""

    for file in uploaded_files:
        if file.type == "text/plain":
            full_text += file.read().decode("utf-8")

        elif file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text

        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = pptx.Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text

    st.sidebar.success("✅ Documents loaded")
    st.session_state.docs = process_documents(full_text)

# ----------------------------------------------------
# CHAT INTERFACE
# ----------------------------------------------------
if "messages" not in st.session_state:
    st.session_state.messages = []

if "docs" in st.session_state:

    st.subheader("💬 Ask Questions")

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    user_input = st.chat_input("Ask something from your documents...")

    if user_input:
        st.session_state.messages.append(
            {"role": "user", "content": user_input}
        )

        with st.chat_message("assistant"):
            with st.spinner("🔍 Searching documents..."):
                result = rag_query(
                    st.session_state.docs,
                    user_input
                )

                st.markdown(f"**Answer:** {result.answer}")
                st.markdown(f"**Confidence:** `{result.confidence}`")

                with st.expander("📄 Source Chunks"):
                    for i, src in enumerate(result.source_chunks, 1):
                        st.markdown(f"**Source {i}:**")
                        st.write(src)

        st.session_state.messages.append(
            {"role": "assistant", "content": result.answer}
        )

# ----------------------------------------------------
# FOOTER
# ----------------------------------------------------
st.markdown("---")
st.markdown("✅ **Groq LLM | Token-safe | Silent embedding fallback | Streamlit**")

